"""
Generate PowerPoint presentation for Lightning Agents talk.

Usage via agent:
    lightning run presentation_slide_writer "List slides"
    lightning run presentation_slide_writer "Add a slide about X"
    lightning run presentation_slide_writer "Generate the presentation"

The agent uses MCP tools: list_slides, add_slide, update_slide, delete_slide, generate_pptx

Outputs both PPTX and PDF (PDF requires LibreOffice, PowerPoint, or Keynote).

Supports simple markup in bullet text:
    **bold text** - renders as bold (primary blue)
    `code text`   - renders as monospace (secondary orange)
"""

import re
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .slide_content import SLIDES, DIAGRAMS
from .styles import COLORS, FONTS, SIZES, DIMS


# Regex patterns for markup parsing
MARKUP_PATTERN = re.compile(r'(\*\*[^*]+\*\*|`[^`]+`)')


def parse_markup(text: str) -> list[tuple[str, str]]:
    """Parse text with **bold** and `code` markup into segments.

    Returns list of (text, style) tuples where style is:
        'normal' - regular text
        'bold'   - bold text
        'code'   - monospace with highlight
    """
    segments = []
    last_end = 0

    for match in MARKUP_PATTERN.finditer(text):
        # Add any text before this match as normal
        if match.start() > last_end:
            segments.append((text[last_end:match.start()], 'normal'))

        matched = match.group()
        if matched.startswith('**') and matched.endswith('**'):
            # Bold text
            segments.append((matched[2:-2], 'bold'))
        elif matched.startswith('`') and matched.endswith('`'):
            # Code text
            segments.append((matched[1:-1], 'code'))

        last_end = match.end()

    # Add remaining text after last match
    if last_end < len(text):
        segments.append((text[last_end:], 'normal'))

    # If no matches, return whole text as normal
    if not segments:
        segments.append((text, 'normal'))

    return segments


def add_rich_text(paragraph, text: str, base_size: int = SIZES["body"]) -> None:
    """Add text with markup to a paragraph using multiple runs."""
    segments = parse_markup(text)

    for i, (segment_text, style) in enumerate(segments):
        if i == 0:
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        else:
            run = paragraph.add_run()

        run.text = segment_text

        if style == 'bold':
            run.font.bold = True
            run.font.color.rgb = rgb("primary")
            run.font.name = FONTS["body"]
            run.font.size = Pt(base_size)
        elif style == 'code':
            run.font.name = FONTS["code"]
            run.font.size = Pt(base_size - 2)
            run.font.color.rgb = rgb("secondary")
        else:
            run.font.name = FONTS["body"]
            run.font.size = Pt(base_size)
            run.font.color.rgb = rgb("text_dark")


def rgb(color_key: str) -> RGBColor:
    """Convert color key to RGBColor."""
    hex_color = COLORS.get(color_key, color_key)
    return RGBColor.from_string(hex_color)


def add_slide_title(slide, title_text: str) -> None:
    """Add a consistent title to any slide."""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4),
        Inches(DIMS["width"] - 1), Inches(1)
    )
    tf = title_box.text_frame
    tf.paragraphs[0].text = title_text
    tf.paragraphs[0].font.size = Pt(SIZES["heading"])
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = rgb("text_dark")
    tf.paragraphs[0].font.name = FONTS["title"]


def add_title_slide(prs: Presentation, data: dict) -> None:
    """Add a centered title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Main title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(DIMS["width"] - 1), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.paragraphs[0].text = data["title"]
    tf.paragraphs[0].font.size = Pt(SIZES["title"])
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = rgb("primary")
    tf.paragraphs[0].font.name = FONTS["title"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    if "subtitle" in data:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2),
            Inches(DIMS["width"] - 1), Inches(1)
        )
        tf = sub_box.text_frame
        tf.paragraphs[0].text = data["subtitle"]
        tf.paragraphs[0].font.size = Pt(SIZES["subtitle"])
        tf.paragraphs[0].font.color.rgb = rgb("text_light")
        tf.paragraphs[0].font.name = FONTS["body"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Presenter name
    if "presenter" in data:
        presenter_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.2),
            Inches(DIMS["width"] - 1), Inches(0.6)
        )
        tf = presenter_box.text_frame
        tf.paragraphs[0].text = data["presenter"]
        tf.paragraphs[0].font.size = Pt(SIZES["body"])
        tf.paragraphs[0].font.color.rgb = rgb("text_dark")
        tf.paragraphs[0].font.name = FONTS["body"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Footer
    if "footer" in data:
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5),
            Inches(DIMS["width"] - 1), Inches(0.5)
        )
        tf = footer_box.text_frame
        tf.paragraphs[0].text = data["footer"]
        tf.paragraphs[0].font.size = Pt(SIZES["small"])
        tf.paragraphs[0].font.color.rgb = rgb("secondary")
        tf.paragraphs[0].font.name = FONTS["body"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_bullet_slide(prs: Presentation, data: dict) -> None:
    """Add a slide with bullet points. Supports **bold** and `code` markup."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_slide_title(slide, data["title"])

    # Bullets
    bullet_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.6),
        Inches(DIMS["width"] - 2), Inches(5.5)
    )
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(data["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(16)
        p.level = 0

        # Add bullet character then rich text
        bullet_run = p.add_run()
        bullet_run.text = "• "
        bullet_run.font.size = Pt(SIZES["body"])
        bullet_run.font.color.rgb = rgb("secondary")

        add_rich_text(p, bullet, SIZES["body"])


def add_code_slide(prs: Presentation, data: dict) -> None:
    """Add a slide with a code block. Optionally includes bullets below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"])

    # Adjust heights based on whether we have bullets
    has_bullets = "bullets" in data and data["bullets"]
    code_box_height = 3.6 if has_bullets else 5.4
    code_text_height = 3.2 if has_bullets else 5.0

    # Code box background
    code_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.6),
        Inches(DIMS["width"] - 1.6), Inches(code_box_height)
    )
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = rgb("code_bg")
    code_bg.line.color.rgb = rgb("text_light")

    # Code text
    code_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.8),
        Inches(DIMS["width"] - 2), Inches(code_text_height)
    )
    tf = code_box.text_frame
    tf.word_wrap = False

    lines = data["code"].split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONTS["code"]
        p.font.size = Pt(SIZES["code"])
        p.font.color.rgb = rgb("text_dark")

    # Optional bullets below code
    if has_bullets:
        bullet_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.4),
            Inches(DIMS["width"] - 2), Inches(1.8)
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(data["bullets"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(8)
            p.level = 0

            # Add bullet character then rich text
            bullet_run = p.add_run()
            bullet_run.text = "• "
            bullet_run.font.size = Pt(SIZES["body"] - 2)
            bullet_run.font.color.rgb = rgb("secondary")

            add_rich_text(p, bullet, SIZES["body"] - 2)


def add_code_comparison_slide(prs: Presentation, data: dict) -> None:
    """Add a slide with side-by-side code comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"])

    half_width = (DIMS["width"] - 1.5) / 2

    # Left side
    left_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(half_width), Inches(0.5)
    )
    tf = left_title.text_frame
    tf.paragraphs[0].text = data["left_title"]
    tf.paragraphs[0].font.size = Pt(SIZES["body"])
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = rgb("text_light")
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    left_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2.0),
        Inches(half_width), Inches(4.5)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = rgb("code_bg")
    left_bg.line.color.rgb = rgb("text_light")

    left_code = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.2),
        Inches(half_width - 0.4), Inches(4.1)
    )
    tf = left_code.text_frame
    for i, line in enumerate(data["left_code"].split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONTS["code"]
        p.font.size = Pt(SIZES["code"])
        p.font.color.rgb = rgb("text_dark")

    # Right side
    right_x = 0.5 + half_width + 0.5

    right_title = slide.shapes.add_textbox(
        Inches(right_x), Inches(1.5),
        Inches(half_width), Inches(0.5)
    )
    tf = right_title.text_frame
    tf.paragraphs[0].text = data["right_title"]
    tf.paragraphs[0].font.size = Pt(SIZES["body"])
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = rgb("secondary")
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    right_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(right_x), Inches(2.0),
        Inches(half_width), Inches(4.5)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = rgb("code_bg")
    right_bg.line.color.rgb = rgb("secondary")

    right_code = slide.shapes.add_textbox(
        Inches(right_x + 0.2), Inches(2.2),
        Inches(half_width - 0.4), Inches(4.1)
    )
    tf = right_code.text_frame
    for i, line in enumerate(data["right_code"].split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONTS["code"]
        p.font.size = Pt(SIZES["code"])
        p.font.color.rgb = rgb("text_dark")


def add_diagram_slide(prs: Presentation, data: dict) -> None:
    """Add a slide with a flow diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"])

    diagram = DIAGRAMS[data["diagram_id"]]

    # Draw boxes
    for box in diagram["boxes"]:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(box["x"]), Inches(box["y"]),
            Inches(2.2), Inches(1.4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(box["color"])
        shape.line.color.rgb = rgb(box["color"])

        # Text in box
        tf = shape.text_frame
        tf.paragraphs[0].text = box["label"]
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = rgb("white")
        tf.paragraphs[0].font.name = FONTS["body"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

    # Draw arrows between boxes
    for start_idx, end_idx in diagram["arrows"]:
        start_box = diagram["boxes"][start_idx]
        end_box = diagram["boxes"][end_idx]

        # Horizontal arrow
        if abs(start_box["y"] - end_box["y"]) < 0.5:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(start_box["x"] + 2.3),
                Inches(start_box["y"] + 0.55),
                Inches(end_box["x"] - start_box["x"] - 2.5),
                Inches(0.3)
            )
        else:
            # Vertical arrow (down)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(start_box["x"] + 0.95),
                Inches(start_box["y"] + 1.5),
                Inches(0.3),
                Inches(end_box["y"] - start_box["y"] - 1.6)
            )

        arrow.fill.solid()
        arrow.fill.fore_color.rgb = rgb("text_light")
        arrow.line.fill.background()


def add_quote_slide(prs: Presentation, data: dict) -> None:
    """Add a slide with a big centered quote and attribution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Quote text (big, centered)
    quote_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.0),
        Inches(DIMS["width"] - 3), Inches(3)
    )
    tf = quote_box.text_frame
    tf.word_wrap = True

    # Opening quote mark
    p = tf.paragraphs[0]
    p.text = f'"{data["quote"]}"'
    p.font.size = Pt(36)
    p.font.italic = True
    p.font.color.rgb = rgb("text_dark")
    p.font.name = FONTS["body"]
    p.alignment = PP_ALIGN.CENTER

    # Attribution
    if "attribution" in data:
        attr_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(5.0),
            Inches(DIMS["width"] - 3), Inches(1)
        )
        tf = attr_box.text_frame
        tf.paragraphs[0].text = f"— {data['attribution']}"
        tf.paragraphs[0].font.size = Pt(SIZES["body"])
        tf.paragraphs[0].font.color.rgb = rgb("secondary")
        tf.paragraphs[0].font.name = FONTS["body"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_comparison_slide(prs: Presentation, data: dict) -> None:
    """Add a slide with a side-by-side comparison table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"])

    half_width = (DIMS["width"] - 1.5) / 2
    header_height = 0.8
    row_height = 0.7
    start_y = 1.8

    # Left column header
    left_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(start_y),
        Inches(half_width), Inches(header_height)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = rgb(data.get("left_color", "text_light"))
    left_header.line.fill.background()
    tf = left_header.text_frame
    tf.paragraphs[0].text = data["left_header"]
    tf.paragraphs[0].font.size = Pt(SIZES["body"])
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = rgb("white")
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    # Right column header
    right_x = 0.5 + half_width + 0.5
    right_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(right_x), Inches(start_y),
        Inches(half_width), Inches(header_height)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = rgb(data.get("right_color", "secondary"))
    right_header.line.fill.background()
    tf = right_header.text_frame
    tf.paragraphs[0].text = data["right_header"]
    tf.paragraphs[0].font.size = Pt(SIZES["body"])
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = rgb("white")
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    # Rows with alternating backgrounds
    row_y = start_y + header_height + 0.1
    for i, (left_text, right_text) in enumerate(data["rows"]):
        # Row background (alternating light gray)
        if i % 2 == 0:
            left_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(row_y),
                Inches(half_width), Inches(row_height)
            )
            left_bg.fill.solid()
            left_bg.fill.fore_color.rgb = rgb("code_bg")
            left_bg.line.fill.background()

            right_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(right_x), Inches(row_y),
                Inches(half_width), Inches(row_height)
            )
            right_bg.fill.solid()
            right_bg.fill.fore_color.rgb = rgb("code_bg")
            right_bg.line.fill.background()

        # Left cell text
        left_cell = slide.shapes.add_textbox(
            Inches(0.5), Inches(row_y + 0.15),
            Inches(half_width), Inches(row_height - 0.3)
        )
        tf = left_cell.text_frame
        tf.paragraphs[0].text = left_text
        tf.paragraphs[0].font.size = Pt(SIZES["body"])
        tf.paragraphs[0].font.color.rgb = rgb("text_dark")
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Right cell text
        right_cell = slide.shapes.add_textbox(
            Inches(right_x), Inches(row_y + 0.15),
            Inches(half_width), Inches(row_height - 0.3)
        )
        tf = right_cell.text_frame
        tf.paragraphs[0].text = right_text
        tf.paragraphs[0].font.size = Pt(SIZES["body"])
        tf.paragraphs[0].font.color.rgb = rgb("text_dark")
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        row_y += row_height

    # Footer tagline
    if "footer" in data:
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.2),
            Inches(DIMS["width"] - 1), Inches(0.8)
        )
        tf = footer_box.text_frame
        tf.paragraphs[0].text = data["footer"]
        tf.paragraphs[0].font.size = Pt(SIZES["body"])
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = rgb("primary")
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_before_after_slide(prs: Presentation, data: dict) -> None:
    """Add a slide with before/after split comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"])

    half_width = (DIMS["width"] - 1.5) / 2
    box_top = 1.8
    box_height = 4.8

    # BEFORE side (left, red accent)
    before_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(box_top),
        Inches(half_width), Inches(box_height)
    )
    before_bg.fill.solid()
    before_bg.fill.fore_color.rgb = rgb("code_bg")
    before_bg.line.color.rgb = rgb("danger")
    before_bg.line.width = Pt(3)

    # Before header
    before_header = slide.shapes.add_textbox(
        Inches(0.5), Inches(box_top + 0.2),
        Inches(half_width), Inches(0.6)
    )
    tf = before_header.text_frame
    tf.paragraphs[0].text = data.get("before_title", "BEFORE")
    tf.paragraphs[0].font.size = Pt(SIZES["body"])
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = rgb("danger")
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Before content
    before_content = slide.shapes.add_textbox(
        Inches(0.7), Inches(box_top + 0.9),
        Inches(half_width - 0.4), Inches(box_height - 1.2)
    )
    tf = before_content.text_frame
    tf.word_wrap = True
    for i, line in enumerate(data["before_items"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"✗ {line}"
        p.font.size = Pt(SIZES["body"] - 2)
        p.font.color.rgb = rgb("text_dark")
        p.space_before = Pt(8)

    # AFTER side (right, green accent)
    right_x = 0.5 + half_width + 0.5
    after_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(right_x), Inches(box_top),
        Inches(half_width), Inches(box_height)
    )
    after_bg.fill.solid()
    after_bg.fill.fore_color.rgb = rgb("code_bg")
    after_bg.line.color.rgb = rgb("success")
    after_bg.line.width = Pt(3)

    # After header
    after_header = slide.shapes.add_textbox(
        Inches(right_x), Inches(box_top + 0.2),
        Inches(half_width), Inches(0.6)
    )
    tf = after_header.text_frame
    tf.paragraphs[0].text = data.get("after_title", "AFTER")
    tf.paragraphs[0].font.size = Pt(SIZES["body"])
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = rgb("success")
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # After content
    after_content = slide.shapes.add_textbox(
        Inches(right_x + 0.2), Inches(box_top + 0.9),
        Inches(half_width - 0.4), Inches(box_height - 1.2)
    )
    tf = after_content.text_frame
    tf.word_wrap = True
    for i, line in enumerate(data["after_items"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"✓ {line}"
        p.font.size = Pt(SIZES["body"] - 2)
        p.font.color.rgb = rgb("text_dark")
        p.space_before = Pt(8)


def add_convergence_slide(prs: Presentation, data: dict) -> None:
    """Add a slide showing multiple paths converging to one point.

    Sources can be:
    - Simple strings: ["Voyager", "LangGraph"]
    - Dicts with optional images: [{"label": "Voyager", "image": "voyager.png"}]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"])

    # Source boxes on the left - compact vertical layout
    sources = data["sources"]
    num_sources = len(sources)
    source_width = 3.0
    source_height = 0.9  # Compact to fit all sources
    source_x = 1.0
    start_y = 1.6
    spacing = 1.15  # Tighter spacing for better arrow alignment

    # Get images directory
    images_dir = Path(__file__).parent / "images"

    source_positions = []
    for i, source in enumerate(sources):
        # Handle both string and dict formats
        if isinstance(source, dict):
            label = source.get("label", "")
            image_file = source.get("image")
        else:
            label = source
            image_file = None

        y = start_y + (i * spacing)

        # Add background box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(source_x), Inches(y),
            Inches(source_width), Inches(source_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = rgb("code_bg")
        box.line.color.rgb = rgb("text_light")

        # Try to add image if specified
        image_path = images_dir / image_file if image_file else None
        if image_path and image_path.exists():
            # Add image on left side of box
            img_size = 0.55
            try:
                slide.shapes.add_picture(
                    str(image_path),
                    Inches(source_x + 0.2), Inches(y + (source_height - img_size) / 2),
                    Inches(img_size), Inches(img_size)
                )
                # Text to the right of image
                text_box = slide.shapes.add_textbox(
                    Inches(source_x + 0.9), Inches(y + 0.1),
                    Inches(source_width - 1.1), Inches(source_height - 0.2)
                )
            except Exception:
                # Fall back to text-only if image fails
                text_box = slide.shapes.add_textbox(
                    Inches(source_x + 0.1), Inches(y + 0.1),
                    Inches(source_width - 0.2), Inches(source_height - 0.2)
                )
        else:
            # Text-only box (centered)
            text_box = slide.shapes.add_textbox(
                Inches(source_x + 0.1), Inches(y + 0.1),
                Inches(source_width - 0.2), Inches(source_height - 0.2)
            )

        tf = text_box.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = rgb("text_dark")
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT if image_path and image_path.exists() else PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Store right edge center point
        source_positions.append((source_x + source_width, y + source_height / 2))

    # Target box on the right - taller to receive all arrows
    target_width = 3.5
    target_height = 2.8  # Taller to span all arrow endpoints
    target_x = 8.5
    # Center target vertically with all sources
    total_sources_height = (num_sources - 1) * spacing + source_height
    middle_y = start_y + total_sources_height / 2
    target_y = middle_y - target_height / 2

    target_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(target_x), Inches(target_y),
        Inches(target_width), Inches(target_height)
    )
    target_box.fill.solid()
    target_box.fill.fore_color.rgb = rgb("primary")
    target_box.line.fill.background()
    tf = target_box.text_frame
    tf.paragraphs[0].text = data["target"]
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = rgb("white")
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    # Draw lines from each source to target left edge center
    target_left_x = target_x
    target_center_y = target_y + target_height / 2

    for src_right_x, src_center_y in source_positions:
        # Draw a line from source right edge to target left edge
        # Using a thin rectangle as a "line"
        line_start_x = src_right_x + 0.1
        line_end_x = target_left_x - 0.1

        # Calculate line angle
        line_width = line_end_x - line_start_x

        # Create arrow shape pointing right
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(line_start_x),
            Inches(src_center_y - 0.1),
            Inches(line_width),
            Inches(0.2)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = rgb("muted")
        arrow.line.fill.background()

    # Footer tagline - positioned below the content
    if "footer" in data:
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.6),
            Inches(DIMS["width"] - 1), Inches(0.6)
        )
        tf = footer_box.text_frame
        tf.paragraphs[0].text = data["footer"]
        tf.paragraphs[0].font.size = Pt(SIZES["body"])
        tf.paragraphs[0].font.italic = True
        tf.paragraphs[0].font.color.rgb = rgb("text_light")
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_closing_slide(prs: Presentation, data: dict) -> None:
    """Add a closing slide with resources."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(DIMS["width"] - 1), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.paragraphs[0].text = data["title"]
    tf.paragraphs[0].font.size = Pt(SIZES["title"])
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = rgb("primary")
    tf.paragraphs[0].font.name = FONTS["title"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Bullets (centered) - supports **bold** and `code` markup
    bullet_box = slide.shapes.add_textbox(
        Inches(2), Inches(3.0),
        Inches(DIMS["width"] - 4), Inches(3)
    )
    tf = bullet_box.text_frame

    for i, bullet in enumerate(data["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(16)
        p.alignment = PP_ALIGN.CENTER
        add_rich_text(p, bullet, SIZES["body"])

    # Footer
    if "footer" in data:
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5),
            Inches(DIMS["width"] - 1), Inches(0.5)
        )
        tf = footer_box.text_frame
        tf.paragraphs[0].text = data["footer"]
        tf.paragraphs[0].font.size = Pt(SIZES["small"])
        tf.paragraphs[0].font.color.rgb = rgb("secondary")
        tf.paragraphs[0].font.name = FONTS["body"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def convert_to_pdf(pptx_path: Path) -> Path | None:
    """Convert PPTX to PDF using available tools.

    Tries in order:
    1. LibreOffice (cross-platform)
    2. Keynote via AppleScript (macOS only)

    Returns PDF path on success, None on failure.
    """
    output_dir = pptx_path.parent
    pdf_path = output_dir / pptx_path.stem.replace(".pptx", "")
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"

    # Try LibreOffice first (works on all platforms)
    soffice_paths = [
        "soffice",  # If in PATH
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",  # macOS
        "/usr/bin/soffice",  # Linux
    ]

    for soffice in soffice_paths:
        try:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(pptx_path)],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode == 0 and pdf_path.exists():
                return pdf_path
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue

    # Try Microsoft PowerPoint on macOS
    if sys.platform == "darwin":
        applescript_ppt = f'''
        tell application "Microsoft PowerPoint"
            open POSIX file "{pptx_path}"
            save active presentation in POSIX file "{pdf_path}" as save as PDF
            close active presentation
        end tell
        '''
        try:
            result = subprocess.run(
                ["osascript", "-e", applescript_ppt],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode == 0 and pdf_path.exists():
                return pdf_path
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # Try Keynote on macOS
        applescript_keynote = f'''
        tell application "Keynote"
            set theDoc to open POSIX file "{pptx_path}"
            export theDoc to POSIX file "{pdf_path}" as PDF
            close theDoc
        end tell
        '''
        try:
            result = subprocess.run(
                ["osascript", "-e", applescript_keynote],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode == 0 and pdf_path.exists():
                return pdf_path
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

    return None


def create_presentation() -> Path:
    """Generate the complete presentation."""
    prs = Presentation()
    prs.slide_width = Inches(DIMS["width"])
    prs.slide_height = Inches(DIMS["height"])

    for slide_data in SLIDES:
        slide_type = slide_data["type"]

        if slide_type == "title":
            add_title_slide(prs, slide_data)
        elif slide_type == "bullets":
            add_bullet_slide(prs, slide_data)
        elif slide_type == "diagram":
            add_diagram_slide(prs, slide_data)
        elif slide_type == "code":
            add_code_slide(prs, slide_data)
        elif slide_type == "code_comparison":
            add_code_comparison_slide(prs, slide_data)
        elif slide_type == "quote":
            add_quote_slide(prs, slide_data)
        elif slide_type == "comparison":
            add_comparison_slide(prs, slide_data)
        elif slide_type == "before_after":
            add_before_after_slide(prs, slide_data)
        elif slide_type == "convergence":
            add_convergence_slide(prs, slide_data)
        elif slide_type == "closing":
            add_closing_slide(prs, slide_data)

    # Save PPTX
    output_dir = Path(__file__).parent / "output"
    output_dir.mkdir(exist_ok=True)
    output_path = output_dir / "lightning-agents.pptx"
    prs.save(output_path)
    print(f"PPTX saved to: {output_path}")

    # Convert to PDF
    pdf_path = convert_to_pdf(output_path)
    if pdf_path:
        print(f"PDF saved to: {pdf_path}")
    else:
        print("PDF conversion skipped (install LibreOffice or use Keynote on macOS)")

    return output_path


def main():
    """CLI entry point."""
    create_presentation()


if __name__ == "__main__":
    main()
